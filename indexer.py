from __future__ import annotations

import json
import hashlib
import math
import os
import re
import threading
import time
import urllib.error
import urllib.request
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import chromadb
from chromadb.config import Settings


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".txt"}
WORD_RE = re.compile(r"[A-Za-z0-9][A-Za-z0-9_+./-]*")
SPEC_INTENT_RE = re.compile(
    r"\b(spec|specs|specification|specifications|dimension|dimensions|capacity|load|rating|part number|sku|model)\b",
    re.IGNORECASE,
)
SPEC_VALUE_RE = re.compile(
    r"\b(\d+(?:[.,]\d+)?\s?(?:U|mm|kg|kgs|kW|W|A|V|Hz|C|F|%|pcs|sets?)|RAL\s?\d{4}|EIA-?310)\b",
    re.IGNORECASE,
)
PARTS_INTENT_RE = re.compile(r"\b(accessory|accessories|part|parts|part number|sku|skus|model|models)\b", re.IGNORECASE)
CORE_SPEC_TERMS = {
    "height", "width", "depth", "static load", "dynamic load", "loading capacity",
    "perforation", "perforated", "eia-310", "eia310", "mounting rails", "standard color",
    "standard colour", "standard inclusions", "side panels", "leveling feet", "casters",
}
ACCESSORY_TABLE_TERMS = {
    "accessories part number", "accessory", "cable manager", "brush panel", "blank panel",
    "fixed shelf", "sliding shelf", "fan tray", "plinth", "lighting panel", "cage nuts",
    "partition panel", "cable ring", "fan top cover",
}
INTENT_TERMS = {
    "compare", "key", "spec", "specs", "specification", "specifications", "dimension", "dimensions",
    "capacity", "load", "rating", "part", "parts", "number", "sku", "skus", "model", "models",
}
STOPWORDS = {
    "a", "about", "an", "and", "are", "as", "at", "be", "by", "can", "do", "does", "for",
    "from", "give", "how", "i", "in", "is", "it", "key", "me", "of", "on", "or", "please",
    "show", "tell", "that", "the", "their", "these", "this", "to", "what", "when", "where",
    "which", "with", "would", "you", "your",
}

QUALITY_TERMS = {
    "accessories", "airflow", "approvals", "capacity", "certifications", "color", "colour",
    "configuration", "depth", "dimension", "dimensions", "door", "dynamic", "eia", "features",
    "frame", "height", "inclusions", "load", "model", "mounting", "number", "part", "pdu",
    "perforation", "rack", "rail", "ral", "regulatory", "sku", "skus", "standard", "static",
    "width",
}
JUNK_TERMS = {
    "all rights reserved", "copyright", "respective owners", "specifications are subject",
    "vertiv.com", "content\n", "additional notes",
}


@dataclass
class Section:
    location: str
    text: str


class HashEmbeddingFunction:
    """Small local embedding function for private, dependency-light Chroma retrieval."""

    def __init__(self, dimensions: int = 384):
        self.dimensions = dimensions

    def name(self) -> str:
        return "vertiv_hash_embeddings"

    def __call__(self, input):
        return [self._embed(text) for text in input]

    def embed_query(self, input):
        return self(input)

    def embed_documents(self, input):
        return self(input)

    def _embed(self, text: str) -> list[float]:
        vector = [0.0] * self.dimensions
        tokens = [token.lower() for token in WORD_RE.findall(text)]
        features = tokens[:700]
        features.extend(f"{a}_{b}" for a, b in zip(tokens, tokens[1:]))
        for token in features:
            digest = hashlib.blake2b(token.encode("utf-8", "ignore"), digest_size=8).digest()
            value = int.from_bytes(digest, "little")
            index = value % self.dimensions
            sign = 1.0 if (value >> 8) & 1 else -1.0
            weight = 1.0 + min(len(token), 24) / 24
            vector[index] += sign * weight
        norm = math.sqrt(sum(value * value for value in vector)) or 1.0
        return [value / norm for value in vector]


class KnowledgeIndex:
    def __init__(self, dataset_dir: Path, vector_path: Path):
        self.dataset_dir = dataset_dir
        self.vector_path = vector_path
        self.manifest_path = vector_path.parent / "chroma_manifest.json"
        self.embedding_function = HashEmbeddingFunction()
        self._lock = threading.Lock()
        self._state_lock = threading.Lock()
        self._status = {
            "running": False,
            "processed": 0,
            "total": 0,
            "current": "",
            "errors": 0,
            "started_at": None,
            "finished_at": None,
        }
        vector_path.mkdir(parents=True, exist_ok=True)
        self.manifest_path.parent.mkdir(parents=True, exist_ok=True)
        self._manifest = self._load_manifest()
        self._initialize()

    def _initialize(self) -> None:
        self.client = chromadb.PersistentClient(
            path=str(self.vector_path),
            settings=Settings(anonymized_telemetry=False),
        )
        self.collection = self.client.get_or_create_collection(
            name="vertiv_documents",
            embedding_function=self.embedding_function,
            metadata={"hnsw:space": "cosine"},
        )

    def _load_manifest(self) -> dict:
        if not self.manifest_path.exists():
            return {"documents": {}}
        try:
            data = json.loads(self.manifest_path.read_text(encoding="utf-8"))
            if isinstance(data, dict) and isinstance(data.get("documents"), dict):
                return data
        except (OSError, json.JSONDecodeError):
            pass
        return {"documents": {}}

    def _save_manifest(self) -> None:
        self.manifest_path.write_text(json.dumps(self._manifest, indent=2), encoding="utf-8")

    def health(self) -> dict:
        return {
            "ok": True,
            "dataset": str(self.dataset_dir),
            "vector_store": str(self.vector_path),
            "retrieval": "ChromaDB",
            "llm": self.llm_config()["enabled"],
        }

    def index_status(self) -> dict:
        with self._state_lock:
            return dict(self._status)

    def stats(self) -> dict:
        manifest_documents = self._manifest.get("documents", {})
        ready_documents = [
            item for item in manifest_documents.values()
            if item.get("status") == "ready"
        ]
        documents = len(ready_documents)
        chunks = self.collection.count()
        errors = sum(1 for item in manifest_documents.values() if item.get("status") == "error")
        category_counts: dict[str, int] = {}
        for item in ready_documents:
            category = item.get("category") or "Other"
            category_counts[category] = category_counts.get(category, 0) + 1
        categories = [
            {"category": category, "count": count}
            for category, count in sorted(category_counts.items(), key=lambda item: item[1], reverse=True)
        ]
        discoverable = sum(1 for p in self.dataset_dir.rglob("*") if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS)
        return {
            "documents": documents,
            "chunks": chunks,
            "errors": errors,
            "discoverable": discoverable,
            "categories": categories,
            "retrieval": "ChromaDB vector search",
            "llm": self.llm_config(),
        }

    def llm_config(self) -> dict:
        model = os.getenv("LLM_MODEL", os.getenv("GROQ_MODEL", "")).strip()
        if not model and os.getenv("GROQ_API_KEY", "").strip():
            model = "llama-3.3-70b-versatile"
        base_url = os.getenv("LLM_BASE_URL", os.getenv("GROQ_BASE_URL", "")).strip()
        if not base_url and os.getenv("GROQ_API_KEY", "").strip():
            base_url = "https://api.groq.com/openai/v1"
        provider = "Groq" if "groq.com" in base_url.lower() else "OpenAI-compatible"
        return {
            "enabled": bool(model),
            "model": model or None,
            "provider": provider if model else "ChromaDB retrieval",
        }

    def start_indexing(self, force: bool = False) -> bool:
        with self._state_lock:
            if self._status["running"]:
                return False
            self._status.update(
                running=True,
                processed=0,
                total=0,
                current="",
                errors=0,
                started_at=time.time(),
                finished_at=None,
            )
        thread = threading.Thread(target=self._index_all, args=(force,), daemon=True)
        thread.start()
        return True

    def _index_all(self, force: bool) -> None:
        files = sorted(
            (p for p in self.dataset_dir.rglob("*") if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS),
            key=lambda p: str(p).lower(),
        )
        with self._state_lock:
            self._status["total"] = len(files)
        seen: set[str] = set()
        try:
            for number, path in enumerate(files, 1):
                relative = path.relative_to(self.dataset_dir).as_posix()
                seen.add(relative)
                with self._state_lock:
                    self._status["current"] = relative
                try:
                    self._index_file(path, relative, force)
                except Exception as exc:
                    self._record_error(path, relative, exc)
                    with self._state_lock:
                        self._status["errors"] += 1
                    print(f"Could not index {relative}: {exc}")
                with self._state_lock:
                    self._status["processed"] = number
            self._remove_missing(seen)
        finally:
            with self._state_lock:
                self._status["running"] = False
                self._status["current"] = ""
                self._status["finished_at"] = time.time()

    def _index_file(self, path: Path, relative: str, force: bool) -> None:
        stat = path.stat()
        existing = self._manifest["documents"].get(relative)
        if (
            not force
            and existing
            and existing.get("size") == stat.st_size
            and existing.get("modified_ns") == stat.st_mtime_ns
            and existing.get("status") in {"ready", "error"}
        ):
            return
        sections = extract_sections(path)
        chunks = []
        for section in sections:
            chunks.extend((section.location, text) for text in chunk_text(section.text))
        if not chunks:
            raise ValueError("No readable text found")
        parts = relative.split("/")
        category = parts[0].replace("_", " ") if len(parts) > 1 else "Other"
        title = clean_title(path.stem)
        document_id = stable_id(relative)
        with self._lock:
            self._delete_document_chunks(relative)
            ids = []
            documents = []
            metadatas = []
            for chunk_index, (location, content) in enumerate(chunks):
                ids.append(f"{document_id}:{chunk_index}")
                documents.append(content)
                metadatas.append(
                    {
                        "document_id": document_id,
                        "path": relative,
                        "title": title,
                        "extension": path.suffix.lower(),
                        "category": category,
                        "location": location,
                        "chunk_index": chunk_index,
                        "url": "/api/file?path=" + __import__("urllib.parse").parse.quote(relative),
                    }
                )
            for start in range(0, len(ids), 128):
                end = start + 128
                self.collection.add(
                    ids=ids[start:end],
                    documents=documents[start:end],
                    metadatas=metadatas[start:end],
                )
            self._manifest["documents"][relative] = {
                "document_id": document_id,
                "title": title,
                "extension": path.suffix.lower(),
                "category": category,
                "size": stat.st_size,
                "modified_ns": stat.st_mtime_ns,
                "status": "ready",
                "error": None,
                "chunks": len(chunks),
                "indexed_at": time.time(),
            }
            self._save_manifest()

    def _record_error(self, path: Path, relative: str, exc: Exception) -> None:
        stat = path.stat()
        category = relative.split("/")[0].replace("_", " ")
        with self._lock:
            self._delete_document_chunks(relative)
            self._manifest["documents"][relative] = {
                "document_id": stable_id(relative),
                "title": clean_title(path.stem),
                "extension": path.suffix.lower(),
                "category": category,
                "size": stat.st_size,
                "modified_ns": stat.st_mtime_ns,
                "status": "error",
                "error": str(exc)[:500],
                "chunks": 0,
                "indexed_at": time.time(),
            }
            self._save_manifest()

    def _remove_missing(self, seen: set[str]) -> None:
        with self._lock:
            for relative in list(self._manifest.get("documents", {})):
                if relative not in seen:
                    self._delete_document_chunks(relative)
                    del self._manifest["documents"][relative]
            self._save_manifest()

    def _delete_document_chunks(self, relative: str) -> None:
        try:
            self.collection.delete(where={"path": relative})
        except Exception:
            pass

    def search(self, query: str, limit: int = 10, category: str = "") -> list[dict]:
        query = query.strip()
        if not query:
            return []
        tokens = [
            token for token in WORD_RE.findall(query)
            if len(token) > 1 and token.lower() not in STOPWORDS
        ][:24]
        if not tokens:
            return []
        return self._chroma_search(query, tokens, limit, category)

    def _chroma_search(self, query: str, tokens: list[str], limit: int, category: str = "") -> list[dict]:
        count = self.collection.count()
        if not count:
            return []
        model_refs = extract_model_references(query)
        query_kwargs = {
            "query_texts": [query],
            "n_results": min(max(limit * 8, 60), count),
            "include": ["documents", "metadatas", "distances"],
        }
        if category:
            query_kwargs["where"] = {"category": category}
        result = self.collection.query(**query_kwargs)
        rows = []
        seen_ids: set[str] = set()
        rows.extend(self._sources_from_query_result(result, tokens, query, seen_ids))
        rows.extend(self._lexical_chroma_sources(query, tokens, category, seen_ids))
        if model_refs:
            for row in rows:
                row["_rank"] += model_match_score(row, model_refs)
            exact_rows = [
                row for row in rows
                if row_matches_all_models(row, model_refs)
            ]
            if exact_rows:
                rows = exact_rows
        rows.sort(key=lambda row: row["_rank"], reverse=True)
        results = []
        per_document: dict[str, int] = {}
        for row in rows:
            document_id = str(row["document_id"])
            if per_document.get(document_id, 0) >= 3:
                continue
            per_document[document_id] = per_document.get(document_id, 0) + 1
            row.pop("_rank", None)
            results.append(row)
            if len(results) >= limit:
                break
        return results

    def _sources_from_query_result(self, result: dict, tokens: list[str], query: str, seen_ids: set[str]) -> list[dict]:
        rows = []
        ids = result.get("ids", [[]])[0] or []
        documents = result.get("documents", [[]])[0] or []
        metadatas = result.get("metadatas", [[]])[0] or []
        distances = result.get("distances", [[]])[0] or []
        for item_id, content, metadata, distance in zip(ids, documents, metadatas, distances):
            if item_id in seen_ids:
                continue
            seen_ids.add(item_id)
            source = self._source_from_chroma(content or "", metadata or {}, float(distance), tokens)
            source["_rank"] = self._rank_search_row(source, query, tokens)
            rows.append(source)
        return rows

    def _lexical_chroma_sources(self, query: str, tokens: list[str], category: str, seen_ids: set[str]) -> list[dict]:
        rows = []
        search_terms = lexical_terms(query, tokens)
        for term in search_terms[:10]:
            try:
                result = self.collection.get(
                    where={"category": category} if category else None,
                    where_document={"$contains": term},
                    limit=180,
                    include=["documents", "metadatas"],
                )
            except Exception:
                continue
            ids = result.get("ids", []) or []
            documents = result.get("documents", []) or []
            metadatas = result.get("metadatas", []) or []
            for item_id, content, metadata in zip(ids, documents, metadatas):
                if item_id in seen_ids:
                    continue
                seen_ids.add(item_id)
                source = self._source_from_chroma(content or "", metadata or {}, 0.35, tokens)
                source["_rank"] = self._rank_search_row(source, query, tokens) + 8
                rows.append(source)
        return rows

    def _source_from_chroma(self, content: str, metadata: dict, distance: float, tokens: list[str]) -> dict:
        path = metadata.get("path", "")
        return {
            "document_id": metadata.get("document_id", ""),
            "title": metadata.get("title", ""),
            "path": path,
            "category": metadata.get("category", ""),
            "location": metadata.get("location", ""),
            "snippet": make_snippet(content, tokens),
            "content": content,
            "score": round(float(distance), 4),
            "url": metadata.get("url") or "/api/file?path=" + __import__("urllib.parse").parse.quote(path),
        }

    def _legacy_sqlite_search_unused(self, query: str, tokens: list[str], limit: int, category: str = "") -> list[dict]:
        return []
        fts_query = " OR ".join(f'"{token.replace(chr(34), "")}"' for token in tokens)
        sql = """
            SELECT document_id, title, path, category, location,
                   snippet(chunks_fts, 0, '<mark>', '</mark>', ' … ', 32) AS snippet,
                   content,
                   bm25(chunks_fts, 1.0, 2.2, 0.8, 0.5) AS score
            FROM chunks_fts WHERE chunks_fts MATCH ?
        """
        parameters: list[object] = [fts_query]
        if category:
            sql += " AND category = ?"
            parameters.append(category)
        sql += " ORDER BY score LIMIT ?"
        parameters.append(max(limit * 8, 60))
        with self.connect() as db:
            rows = db.execute(sql, parameters).fetchall()
        rows = sorted(rows, key=lambda row: self._rank_search_row(row, query, tokens), reverse=True)
        results = []
        per_document: dict[int, int] = {}
        for row in rows:
            document_id = int(row["document_id"])
            if per_document.get(document_id, 0) >= 3:
                continue
            per_document[document_id] = per_document.get(document_id, 0) + 1
            results.append(
                {
                    "document_id": document_id,
                    "title": row["title"],
                    "path": row["path"],
                    "category": row["category"],
                    "location": row["location"],
                    "snippet": row["snippet"],
                    "content": row["content"],
                    "score": round(float(row["score"]), 4),
                    "url": "/api/file?path=" + __import__("urllib.parse").parse.quote(row["path"]),
                }
            )
            if len(results) >= limit:
                break
        return results

    def _rank_search_row(self, row: dict, query: str, tokens: list[str]) -> float:
        content = row["content"] or ""
        title = row["title"] or ""
        haystack_raw = f"{title} {row['path']} {row['location']} {content}"
        haystack = haystack_raw.lower()
        important_tokens = [
            token for token in tokens
            if token.lower() not in INTENT_TERMS and token.lower() not in STOPWORDS
        ]
        token_hits = sum(1 for token in tokens if token_present(haystack_raw, token))
        important_hits = sum(1 for token in important_tokens if token_present(haystack_raw, token))
        title_hits = sum(1 for token in tokens if token_present(title, token))
        quality_hits = sum(1 for term in QUALITY_TERMS if term in haystack)
        value_hits = len(SPEC_VALUE_RE.findall(content))
        junk_hits = sum(1 for term in JUNK_TERMS if term in haystack)
        core_hits = sum(1 for term in CORE_SPEC_TERMS if term in haystack)
        accessory_hits = sum(1 for term in ACCESSORY_TABLE_TERMS if term in haystack)

        score = token_hits * 6 + important_hits * 12 + title_hits * 8 + min(quality_hits, 8) * 2 + min(value_hits, 10)
        if important_tokens:
            missing = len(important_tokens) - important_hits
            score -= missing * 18
            if missing == 0:
                score += 40
        if SPEC_INTENT_RE.search(query):
            score += min(value_hits, 12) + core_hits * 8
            if "specification" in haystack or "specifications" in haystack:
                score += 8
            if "|" in content:
                score += 3
            if accessory_hits and not PARTS_INTENT_RE.search(query):
                score -= min(accessory_hits, 8) * 8
        if len(content.strip()) < 90:
            score -= 8
        if junk_hits:
            score -= junk_hits * 10
        return score - float(row["score"]) * 4

    def answer(self, question: str, history: list, category: str = "") -> dict:
        safe_history = history if isinstance(history, list) else []
        retrieval_query = contextualize_followup_question(question, safe_history)
        sources = self.search(retrieval_query, 8, category)
        if not sources:
            return {
                "answer": "I couldn’t find relevant evidence in the indexed Vertiv documents. Try a model number, product family, or a more specific technical term.",
                "sources": [],
                "mode": "search",
            }
        config = self.llm_config()
        if config["enabled"]:
            try:
                answer = self._llm_answer(question, safe_history, sources, retrieval_query)
                return {"answer": answer, "sources": sources, "mode": "llm"}
            except Exception as exc:
                print(f"LLM request failed: {exc}")
        answer = self._evidence_answer(retrieval_query, sources)
        return {"answer": answer, "sources": sources, "mode": "search"}

    def _evidence_answer(self, question: str, sources: list[dict]) -> str:
        bullets = extract_answer_bullets(question, sources)
        if bullets:
            lines = ["Here is the strongest answer I found in the Vertiv documents:", ""]
            lines.extend(f"- {bullet}" for bullet in bullets[:6])
            lines.extend(["", "Sources are listed below, ordered by relevance."])
            return "\n".join(lines)

        lines = ["I found relevant Vertiv document excerpts, but not enough clean detail to synthesize a confident answer:", ""]
        for number, source in enumerate(sources[:4], 1):
            text = summarize_excerpt(source.get("content") or source["snippet"])
            lines.append(f"- {text} [{number}]")
        return "\n".join(lines)

    def _llm_answer(self, question: str, history: list, sources: list[dict], retrieval_query: str = "") -> str:
        groq_key = os.getenv("GROQ_API_KEY", "").strip()
        base_url = os.getenv("LLM_BASE_URL", os.getenv("GROQ_BASE_URL", "")).strip()
        if not base_url:
            base_url = "https://api.groq.com/openai/v1" if groq_key else "http://127.0.0.1:11434/v1"
        base_url = base_url.rstrip("/")
        api_key = os.getenv("LLM_API_KEY", groq_key or "ollama")
        model = os.getenv("LLM_MODEL", os.getenv("GROQ_MODEL", "")).strip()
        if not model and groq_key:
            model = "llama-3.3-70b-versatile"
        context_parts = []
        for number, source in enumerate(sources, 1):
            excerpt = llm_context_excerpt(source.get("content") or source.get("snippet", ""))
            context_parts.append(
                f"[{number}] {source['title']} — {source['location']}\nPath: {source['path']}\n{excerpt}"
            )
        model_refs = extract_model_references(retrieval_query or question)
        model_instruction = ""
        if model_refs:
            requested = ", ".join(model_ref["code"] for model_ref in model_refs)
            variants = "; ".join(
                f"{model_ref['code']} may appear as {', '.join(model_ref['variants'])}"
                for model_ref in model_refs
            )
            model_instruction = (
                f" The user requested these exact model identifiers: {requested}. "
                f"Use only facts tied to those exact models. Accepted identifier variants: {variants}. "
                "If an excerpt also lists sibling models or nearby capacities, do not copy values from those other columns/rows. "
                "If the exact model value is unclear in a table, say the evidence is insufficient for that field."
            )
        system = (
            "You are Vertiv Knowledge, an internal product-document assistant. Answer only from the supplied excerpts. "
            "Be precise, preserve units and model names, and say when evidence is insufficient. Cite every factual claim "
            "with bracketed source numbers like [1]. Never invent specifications or citations."
            + model_instruction
        )
        messages = [{"role": "system", "content": system}]
        for item in history[-6:]:
            role = item.get("role")
            content = str(item.get("content", ""))[:3000]
            if role in {"user", "assistant"} and content:
                messages.append({"role": role, "content": content})
        messages.append(
            {
                "role": "user",
                "content": (
                    f"Current user question: {question}\n"
                    f"Retrieval query used for document search: {retrieval_query or question}\n\n"
                    "Answer the current user question. Use the retrieval query only to resolve follow-up context; "
                    "do not mention it unless needed for clarity.\n\n"
                    "Document excerpts:\n\n"
                    + "\n\n".join(context_parts)
                ),
            }
        )
        body = json.dumps({"model": model, "messages": messages, "temperature": 0.1, "stream": False}).encode()
        request = urllib.request.Request(
            base_url + "/chat/completions",
            data=body,
            headers={
                "Content-Type": "application/json",
                "Authorization": f"Bearer {api_key}",
                "User-Agent": "VertivKnowledge/1.0",
            },
            method="POST",
        )
        try:
            with urllib.request.urlopen(request, timeout=120) as response:
                result = json.loads(response.read())
        except urllib.error.HTTPError as exc:
            detail = exc.read().decode("utf-8", "replace")[:500]
            raise RuntimeError(f"LLM returned {exc.code}: {detail}") from exc
        return result["choices"][0]["message"]["content"].strip()


def clean_title(stem: str) -> str:
    return re.sub(r"[_-]+", " ", stem).strip()


def stable_id(value: str) -> str:
    return hashlib.sha1(value.encode("utf-8", "ignore")).hexdigest()


def make_snippet(content: str, tokens: list[str], size: int = 420) -> str:
    clean = normalize_text(content)
    if not clean:
        return ""
    lowered = clean.lower()
    positions = [lowered.find(token.lower()) for token in tokens if lowered.find(token.lower()) >= 0]
    if positions:
        start = max(min(positions) - size // 3, 0)
    else:
        start = 0
    end = min(start + size, len(clean))
    snippet = clean[start:end].strip()
    if start > 0:
        snippet = "... " + snippet
    if end < len(clean):
        snippet += " ..."
    for token in sorted(set(tokens), key=len, reverse=True):
        snippet = re.sub(
            rf"({re.escape(token)})",
            r"<mark>\1</mark>",
            snippet,
            flags=re.IGNORECASE,
        )
    return snippet


def lexical_terms(query: str, tokens: list[str]) -> list[str]:
    terms: list[str] = []
    significant = [token for token in tokens if token.lower() not in STOPWORDS and len(token) > 1]
    for token in significant:
        variants = {token}
        if token.islower():
            variants.add(token.upper())
            variants.add(token.title())
        for variant in variants:
            if variant not in terms:
                terms.append(variant)
    for first, second in zip(significant, significant[1:]):
        phrase = f"{first} {second}"
        if phrase not in terms:
            terms.insert(0, phrase)
    model_like = [
        token for token in significant
        if re.search(r"[A-Za-z]+\d+|\d+[A-Za-z]+", token)
    ]
    for token in model_like:
        if token.upper() not in terms:
            terms.insert(0, token.upper())
    for model_ref in extract_model_references(query):
        for variant in model_ref["variants"]:
            if variant not in terms:
                terms.insert(0, variant)
    return terms


def token_present(text: str, token: str) -> bool:
    return bool(re.search(rf"(?<![A-Za-z0-9]){re.escape(token)}(?![A-Za-z0-9])", text, re.IGNORECASE))


def compact_code(value: str) -> str:
    return re.sub(r"[^A-Za-z0-9]", "", value).upper()


def extract_model_references(text: str) -> list[dict]:
    """Extract exact product/model identifiers from user text.

    This intentionally targets mixed alpha/numeric identifiers (CRV036, CR036RA,
    P1035UA, DME27) and ignores generic family-only terms such as "CRV" or
    intent-only numbers such as "details".
    """
    refs: list[dict] = []
    seen: set[str] = set()
    raw_tokens = WORD_RE.findall(text)
    candidates: list[str] = []
    candidates.extend(raw_tokens)
    for first, second in zip(raw_tokens, raw_tokens[1:]):
        if re.fullmatch(r"[A-Za-z]{2,8}", first) and re.fullmatch(r"\d{2,5}[A-Za-z0-9]*", second):
            candidates.append(first + second)
    for candidate in candidates:
        code = compact_code(candidate)
        if not re.fullmatch(r"[A-Z]{1,10}\d{2,6}[A-Z0-9]*", code):
            continue
        if code in seen:
            continue
        match = re.match(r"^([A-Z]+)(\d{2,6})([A-Z0-9]*)$", code)
        if not match:
            continue
        prefix, number, suffix = match.groups()
        variants = {code}
        base = prefix + number
        variants.add(base)
        if suffix:
            variants.add(prefix + number + suffix)
        # Common Vertiv cooling shorthand: users say CRV036, while rating
        # sheets often encode the same unit as CR036RA. Keep this generic and
        # conservative: only derive it when the alpha prefix ends in V and has
        # at least three letters.
        if prefix.endswith("V") and len(prefix) >= 3:
            variants.add(prefix[:-1] + number)
        refs.append(
            {
                "code": code,
                "prefix": prefix,
                "number": number,
                "variants": sorted(variants, key=len, reverse=True),
            }
        )
        seen.add(code)
    return refs


def recent_history_model_references(history: list, max_refs: int = 3) -> list[dict]:
    refs: list[dict] = []
    seen: set[str] = set()
    if not isinstance(history, list):
        return refs

    # Prefer the user's own recent turns over assistant answers, because
    # assistant answers may contain comparison/sibling model codes from tables.
    for preferred_role in ("user", "assistant"):
        for item in reversed(history[-10:]):
            if not isinstance(item, dict) or item.get("role") != preferred_role:
                continue
            content = str(item.get("content", ""))[:3000]
            for model_ref in extract_model_references(content):
                code = model_ref["code"]
                if code in seen:
                    continue
                refs.append(model_ref)
                seen.add(code)
                if len(refs) >= max_refs:
                    return refs
        if refs:
            return refs
    return refs


def contextualize_followup_question(question: str, history: list) -> str:
    if extract_model_references(question):
        return question
    history_refs = recent_history_model_references(history)
    if not history_refs:
        return question
    model_context = " ".join(model_ref["code"] for model_ref in history_refs)
    return f"{question} {model_context}".strip()


def normalized_code_tokens(text: str) -> list[str]:
    return [compact_code(token) for token in WORD_RE.findall(text) if compact_code(token)]


def code_variant_present(tokens: list[str], variant: str) -> bool:
    for token in tokens:
        if token == variant:
            return True
        if token.startswith(variant):
            remainder = token[len(variant):]
            # CR036 should match CR036RA, but not CR0365.
            if remainder and not remainder[0].isdigit():
                return True
    return False


def row_code_tokens(row: dict) -> list[str]:
    return normalized_code_tokens(
        f"{row.get('title', '')} {row.get('path', '')} {row.get('location', '')} {row.get('content', '')}"
    )


def row_matches_model(row: dict, model_ref: dict) -> bool:
    tokens = row_code_tokens(row)
    return any(code_variant_present(tokens, variant) for variant in model_ref["variants"])


def row_matches_all_models(row: dict, model_refs: list[dict]) -> bool:
    return all(row_matches_model(row, model_ref) for model_ref in model_refs)


def model_match_score(row: dict, model_refs: list[dict]) -> int:
    tokens = row_code_tokens(row)
    score = 0
    for model_ref in model_refs:
        exact = any(code_variant_present(tokens, variant) for variant in model_ref["variants"])
        if exact:
            score += 140
            title_path_tokens = normalized_code_tokens(f"{row.get('title', '')} {row.get('path', '')}")
            if any(code_variant_present(title_path_tokens, variant) for variant in model_ref["variants"]):
                score += 80
            continue
        # Penalize nearby sibling models with same alphabetic prefix but a
        # different numeric capacity/model code. This keeps CR035/CR045 from
        # contaminating a CRV036 answer without hard-coding CRV.
        sibling = False
        requested_prefixes = {model_ref["prefix"]}
        if model_ref["prefix"].endswith("V") and len(model_ref["prefix"]) >= 3:
            requested_prefixes.add(model_ref["prefix"][:-1])
        for token in tokens:
            match = re.match(r"^([A-Z]+)(\d{2,6})", token)
            if match and match.group(1) in requested_prefixes and match.group(2) != model_ref["number"]:
                sibling = True
                break
        score -= 160 if sibling else 80
    return score


def extract_answer_bullets(question: str, sources: list[dict]) -> list[str]:
    query_tokens = {
        token.lower()
        for token in WORD_RE.findall(question)
        if len(token) > 1 and token.lower() not in STOPWORDS
    }
    spec_intent = bool(SPEC_INTENT_RE.search(question))
    parts_intent = bool(PARTS_INTENT_RE.search(question))
    candidates: list[tuple[float, str, str]] = []
    seen: set[str] = set()

    for source_number, source in enumerate(sources[:8], 1):
        text = re.sub(r"</?mark>", "", source.get("content") or source.get("snippet", ""))
        for line in evidence_lines(text):
            normalized = normalize_candidate(line)
            if not normalized or normalized in seen:
                continue
            lower = normalized.lower()
            if any(term.strip() and term in lower for term in JUNK_TERMS):
                continue
            accessory_line = any(term in lower for term in ACCESSORY_TABLE_TERMS)
            if spec_intent and accessory_line and not parts_intent:
                continue
            has_query = any(token_present(normalized, token) for token in query_tokens)
            has_quality = any(term in lower for term in QUALITY_TERMS)
            core_hits = sum(1 for term in CORE_SPEC_TERMS if term in lower)
            value_hits = len(SPEC_VALUE_RE.findall(normalized))
            if spec_intent and not (value_hits or has_quality or has_query):
                continue
            if not spec_intent and not (has_query or value_hits):
                continue

            score = value_hits * 5 + core_hits * 8 + int(has_query) * 4 + int(has_quality) * 3
            if "|" in normalized:
                score += 2
            if len(normalized) > 220:
                score -= 4
            if len(normalized) < 8:
                score -= 3

            seen.add(normalized)
            candidates.append((score, evidence_facet(normalized), f"{humanize_evidence_line(normalized)} [{source_number}]"))

    candidates.sort(key=lambda item: item[0], reverse=True)
    if not spec_intent:
        return [text for _, _, text in candidates]

    selected: list[str] = []
    used_text: set[str] = set()
    facet_order = ["height", "dimensions", "load", "cooling", "rails", "color", "inclusions", "features", "other"]
    for facet in facet_order:
        for _, candidate_facet, text in candidates:
            text_key = re.sub(r"\[[0-9]+\]$", "", text).strip().lower()
            if candidate_facet == facet and text_key not in used_text:
                selected.append(text)
                used_text.add(text_key)
                break
        if len(selected) >= 8:
            break
    return selected or [text for _, _, text in candidates]


def evidence_lines(text: str) -> list[str]:
    text = normalize_text(text)
    raw_lines = []
    for paragraph in re.split(r"[\n\r]+", text):
        paragraph = paragraph.strip()
        if not paragraph:
            continue
        raw_lines.extend(re.split(r"\s+[y•]\s+| yy| • ", paragraph))
    cleaned = [line.strip(" -*\t") for line in raw_lines if line.strip(" -*\t")]
    combined = []
    index = 0
    while index < len(cleaned):
        current = cleaned[index]
        if index + 1 < len(cleaned) and current.lower().endswith(" and") and len(cleaned[index + 1]) <= 40:
            current = f"{current} {cleaned[index + 1]}"
            index += 1
        combined.append(current)
        index += 1
    return combined


def normalize_candidate(line: str) -> str:
    line = re.sub(r"\s+", " ", line)
    line = re.sub(r"\s*\|\s*", " | ", line)
    line = re.sub(r"^(y+|â€¢|\u2022)\s*", "", line, flags=re.IGNORECASE)
    line = line.replace("kgs", "kg").strip(" |")
    if line.lower() in {"content", "vertiv ve rack", "specifications", "configurations", "additional notes", "standard inclusions"}:
        return ""
    if len(line) > 320:
        return ""
    return line.strip()


def humanize_evidence_line(line: str) -> str:
    line = re.sub(r"\b(\d+)\s?kgs\b", r"\1 kg", line, flags=re.IGNORECASE)
    line = re.sub(r"\bfeets\b", "feet", line, flags=re.IGNORECASE)
    return line[:1].upper() + line[1:]


def evidence_facet(line: str) -> str:
    lower = line.lower()
    if "static load" in lower or "dynamic load" in lower or "loading capacity" in lower:
        return "load"
    if "42u" in lower or "48u" in lower or "enclosure height" in lower or lower.startswith("height "):
        return "height"
    if "width" in lower or "depth" in lower or re.search(r"\b\d+\s?x\s?\d+", lower):
        return "dimensions"
    if "perforation" in lower or "perforated" in lower or "air flow" in lower or "airflow" in lower or "cooling" in lower:
        return "cooling"
    if "eia" in lower or "mounting rail" in lower:
        return "rails"
    if "ral" in lower or "color" in lower or "colour" in lower:
        return "color"
    if "standard inclusion" in lower or "included" in lower:
        return "inclusions"
    if any(term in lower for term in ("side panel", "leveling feet", "casters", "pdu", "baying")):
        return "features"
    return "other"


def summarize_excerpt(text: str) -> str:
    text = re.sub(r"</?mark>", "", text)
    lines = evidence_lines(text)
    if not lines:
        return normalize_text(text)[:240]
    return " ".join(lines[:3])[:260]


def llm_context_excerpt(text: str, limit: int = 2400) -> str:
    text = re.sub(r"</?mark>", "", text)
    text = normalize_text(text)
    if len(text) <= limit:
        return text
    return text[:limit].rstrip() + " ..."


def normalize_text(text: str) -> str:
    text = text.replace("\x00", " ").replace("\u00ad", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def chunk_text(text: str, target: int = 1400, overlap: int = 220) -> Iterable[str]:
    text = normalize_text(text)
    if not text:
        return
    start = 0
    while start < len(text):
        end = min(start + target, len(text))
        if end < len(text):
            boundary = max(text.rfind("\n", start + target // 2, end), text.rfind(". ", start + target // 2, end))
            if boundary > start:
                end = boundary + 1
        chunk = text[start:end].strip()
        if len(chunk) >= 40:
            yield chunk
        if end >= len(text):
            break
        start = max(end - overlap, start + 1)


def extract_sections(path: Path) -> list[Section]:
    extension = path.suffix.lower()
    if extension == ".pdf":
        from pypdf import PdfReader

        reader = PdfReader(str(path), strict=False)
        return [Section(f"Page {number}", page.extract_text() or "") for number, page in enumerate(reader.pages, 1)]
    if extension == ".docx":
        from docx import Document

        document = Document(str(path))
        text = [paragraph.text for paragraph in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                text.append(" | ".join(cell.text for cell in row.cells))
        return [Section("Document", "\n".join(text))]
    if extension == ".pptx":
        from pptx import Presentation

        presentation = Presentation(str(path))
        sections = []
        for number, slide in enumerate(presentation.slides, 1):
            text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        text.append(" | ".join(cell.text for cell in row.cells))
            sections.append(Section(f"Slide {number}", "\n".join(text)))
        return sections
    if extension == ".xlsx":
        from openpyxl import load_workbook

        workbook = load_workbook(path, read_only=True, data_only=True)
        sections = []
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(value) if value is not None else "" for value in row]
                if any(values):
                    rows.append(" | ".join(values))
            sections.append(Section(f"Sheet: {sheet.title}", "\n".join(rows)))
        workbook.close()
        return sections
    if extension == ".txt":
        return [Section("Text file", path.read_text(encoding="utf-8", errors="replace"))]
    raise ValueError(f"Unsupported format: {extension}")
